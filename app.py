import io
import os
import subprocess
import tempfile
import uuid
import zipfile
from pathlib import Path

from flask import (
    Flask,
    render_template,
    request,
    send_file,
    jsonify,
    after_this_request,
    Response,
)
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
from PIL import Image, ImageOps
from PIL import UnidentifiedImageError
from pypdf import PdfReader, PdfWriter
from pypdf.errors import PdfReadError
from pdf2docx import Converter
import fitz  # PyMuPDF
import qrcode
import pytesseract
from PIL import Image as PILImage
import pillow_heif
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.colors import Color
from pptx import Presentation
from pptx.util import Emu
import pdfplumber
from openpyxl import Workbook

pillow_heif.register_heif_opener()

app = Flask(__name__)

SITE_NAME = "사무실 공구함"
SITE_TAGLINE = "OFFICE TOOLBOX"

MAX_CONTENT_LENGTH = 100 * 1024 * 1024  # 100MB (PDF 쪽이 더 큼)
app.config["MAX_CONTENT_LENGTH"] = MAX_CONTENT_LENGTH

PDF_MAX_FILE_BYTES = 100 * 1024 * 1024
IMAGE_COMPRESS_MAX_FILE_BYTES = 30 * 1024 * 1024
IMAGE_CONVERT_MAX_FILE_BYTES = 50 * 1024 * 1024
OCR_MAX_FILE_BYTES = 50 * 1024 * 1024
SIGNATURE_MAX_FILE_BYTES = 5 * 1024 * 1024

# Render의 제한된 메모리에서 확인한 안전 범위. 파일 용량이 작아도 지나치게
# 큰 해상도는 디코딩·OCR 과정에서 워커를 종료시킬 수 있으므로 별도로 제한한다.
IMAGE_COMPRESS_MAX_PIXELS = 40_000_000
IMAGE_CONVERT_MAX_PIXELS = 16_000_000
OCR_IMAGE_MAX_PIXELS = 16_000_000
OCR_PDF_MAX_PAGE_PIXELS = 16_000_000
OCR_PAGE_TIMEOUT_SECONDS = 45
SIGNATURE_MAX_PIXELS = 4_000_000

# 사이트맵에 반영하는 마지막 콘텐츠 갱신일. 페이지 구성/문구를 의미 있게
# 바꿀 때마다 이 값을 갱신해 검색엔진에 재수집 신호를 준다.
SITE_LAST_UPDATED = "2026-09-02"


@app.context_processor
def inject_analytics():
    # Render 환경변수 GA_MEASUREMENT_ID가 설정된 경우에만 GA4 스크립트를 렌더링한다.
    return {"ga_measurement_id": os.environ.get("GA_MEASUREMENT_ID", "")}


# ---------------------------------------------------------------------------
# 툴 목록 (홈 화면 카드 + sitemap 생성에 공용으로 사용)
# ---------------------------------------------------------------------------
TOOLS = [
    {
        "slug": "pdf-compress",
        "icon": "pdf-compress",
        "title": "PDF 압축",
        "desc": "용량 큰 PDF를 화질 손상 최소화하며 빠르게 줄여요.",
        "available": True,
        "category": "pdf",
        "popular": True,
    },
    {
        "slug": "image-compress",
        "icon": "image-compress",
        "title": "이미지 압축",
        "desc": "JPG·PNG·WEBP 이미지를 리사이즈하고 압축해요.",
        "available": True,
        "category": "image",
        "popular": True,
    },
    {
        "slug": "pdf-to-word",
        "icon": "pdf-to-word",
        "title": "PDF → Word 변환",
        "desc": "PDF 문서를 편집 가능한 Word(.docx) 파일로 변환해요.",
        "available": True,
        "category": "pdf",
        "popular": True,
    },
    {
        "slug": "pdf-merge-split",
        "icon": "pdf-merge-split",
        "title": "PDF 병합·분할",
        "desc": "여러 PDF를 하나로 합치거나, 한 PDF를 원하는 대로 나눠요.",
        "available": True,
        "category": "pdf",
        "popular": True,
    },
    {
        "slug": "watermark",
        "icon": "watermark",
        "title": "워터마크 추가",
        "desc": "PDF에 원하는 텍스트 워터마크를 삽입해요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "pdf-rotate",
        "icon": "pdf-rotate",
        "title": "PDF 페이지 회전",
        "desc": "PDF 페이지를 원하는 방향으로 회전해요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "pdf-password",
        "icon": "pdf-password",
        "title": "PDF 비밀번호 설정·해제",
        "desc": "PDF에 암호를 걸거나, 알고 있는 암호를 풀어드려요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "char-counter",
        "icon": "char-counter",
        "title": "글자 수·바이트 계산기",
        "desc": "공백 포함/제외 글자수, 바이트, 원고지 매수까지 한 번에 세어드려요.",
        "available": True,
        "category": "text",
        "popular": True,
    },
    {
        "slug": "qr-code",
        "icon": "qr-code",
        "title": "QR코드 생성기",
        "desc": "텍스트나 URL을 입력하면 QR코드 이미지를 만들어드려요.",
        "available": True,
        "category": "text",
        "popular": False,
    },
    {
        "slug": "text-diff",
        "icon": "text-diff",
        "title": "텍스트 비교",
        "desc": "두 텍스트를 비교해서 달라진 부분을 한눈에 보여줘요.",
        "available": True,
        "category": "text",
        "popular": False,
    },
    {
        "slug": "salary-calculator",
        "icon": "salary-calculator",
        "title": "연봉 실수령액 계산기",
        "desc": "4대보험료와 세금을 반영한 월 실수령액을 계산해요.",
        "available": True,
        "category": "calculator",
        "popular": False,
    },
    {
        "slug": "insurance-calculator",
        "icon": "document",
        "title": "4대보험료 계산기",
        "desc": "국민연금·건강보험·고용보험·산재보험료를 근로자·사업주 부담분으로 나눠 계산해요.",
        "available": True,
        "category": "calculator",
        "popular": False,
    },
    {
        "slug": "severance-calculator",
        "icon": "severance-calculator",
        "title": "퇴직금 계산기",
        "desc": "입사일·퇴사일과 급여로 예상 퇴직금을 계산해요.",
        "available": True,
        "category": "calculator",
        "popular": False,
    },
    {
        "slug": "jeonse-rent-calculator",
        "icon": "document",
        "title": "전월세 전환 계산기",
        "desc": "전세보증금을 월세로 전환할 때 적정 월세와 법정 전환율 상한을 계산해요.",
        "available": True,
        "category": "calculator",
        "popular": False,
    },
    {
        "slug": "image-convert",
        "icon": "image-convert",
        "title": "이미지 포맷 변환",
        "desc": "HEIC·JPG·PNG·WEBP 등 이미지 포맷을 서로 변환해요.",
        "available": True,
        "category": "image",
        "popular": False,
    },
    {
        "slug": "pdf-sign",
        "icon": "pdf-sign",
        "title": "전자서명 삽입",
        "desc": "직접 그린 서명을 PDF 원하는 위치에 넣어드려요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "ocr",
        "icon": "ocr",
        "title": "OCR 텍스트 추출",
        "desc": "스캔한 PDF·이미지에서 텍스트를 인식해 추출해요.",
        "available": True,
        "category": "text",
        "popular": False,
    },
    {
        "slug": "pdf-to-ppt",
        "icon": "pdf-to-ppt",
        "title": "PDF → PPT 변환",
        "desc": "PDF 페이지를 슬라이드 이미지로 담은 PPT로 만들어요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "pdf-to-excel",
        "icon": "pdf-to-excel",
        "title": "PDF → Excel 변환",
        "desc": "PDF 안의 표를 추출해 엑셀 파일로 만들어요.",
        "available": True,
        "category": "pdf",
        "popular": False,
    },
    {
        "slug": "quote-statement",
        "icon": "document",
        "title": "견적서·거래명세서 만들기",
        "desc": "품목과 부가세를 계산해 A4 문서를 만들고 PDF로 저장해요.",
        "available": True,
        "category": "business",
        "popular": True,
    },
]

TOOL_CATEGORIES = [
    {"key": "pdf", "label": "PDF", "eyebrow": "DOCUMENT"},
    {"key": "image", "label": "이미지", "eyebrow": "IMAGE"},
    {"key": "text", "label": "텍스트·데이터", "eyebrow": "TEXT & DATA"},
    {"key": "calculator", "label": "업무 계산기", "eyebrow": "CALCULATOR"},
]

# ---------------------------------------------------------------------------
# 업로드 임시 디렉토리
# ---------------------------------------------------------------------------
UPLOAD_DIR = Path(tempfile.gettempdir()) / "office-toolbox-uploads"
UPLOAD_DIR.mkdir(exist_ok=True)


class ClientFileError(ValueError):
    """사용자가 다른 파일을 선택하면 해결되는 입력 오류."""


class UploadTooLargeError(ClientFileError):
    """도구별 업로드 허용 용량을 초과한 경우."""


def save_upload_limited(file_storage, destination: Path, max_bytes: int, label: str) -> int:
    """업로드를 청크 단위로 저장하며 도구별 실제 파일 크기를 강제한다."""
    written = 0
    try:
        with destination.open("wb") as output:
            while True:
                chunk = file_storage.stream.read(1024 * 1024)
                if not chunk:
                    break
                written += len(chunk)
                if written > max_bytes:
                    raise UploadTooLargeError(
                        f"{label} 파일은 최대 {human_size(max_bytes)}까지 처리할 수 있습니다."
                    )
                output.write(chunk)
    except Exception:
        destination.unlink(missing_ok=True)
        raise

    if written == 0:
        destination.unlink(missing_ok=True)
        raise ClientFileError("빈 파일은 처리할 수 없습니다.")
    return written


def validate_pdf_file(path: Path, *, allow_encrypted: bool = False) -> tuple[int, bool]:
    """확장자가 아니라 실제 PDF 구조를 확인한다."""
    try:
        reader = PdfReader(str(path), strict=False)
        encrypted = bool(reader.is_encrypted)
        if encrypted and not allow_encrypted:
            raise ClientFileError("암호가 걸린 PDF는 이 도구에서 처리할 수 없습니다.")
        if encrypted:
            return 0, True
        page_count = len(reader.pages)
        if page_count < 1:
            raise ClientFileError("페이지가 없는 PDF는 처리할 수 없습니다.")
        return page_count, False
    except ClientFileError:
        raise
    except (PdfReadError, EOFError, ValueError, OSError) as exc:
        raise ClientFileError("손상되었거나 올바른 PDF 파일이 아닙니다.") from exc


def validate_image_file(path: Path, *, max_pixels: int) -> tuple[int, int]:
    """이미지 구조와 해상도를 실제 디코딩 전에 확인한다."""
    try:
        with PILImage.open(path) as image:
            width, height = image.size
            if width < 1 or height < 1:
                raise ClientFileError("크기가 올바르지 않은 이미지입니다.")
            pixels = width * height
            if pixels > max_pixels:
                max_mp = max_pixels // 1_000_000
                raise ClientFileError(
                    f"이미지 해상도가 너무 큽니다. {max_mp}메가픽셀 이하 이미지를 사용해주세요."
                )
            image.verify()
        return width, height
    except ClientFileError:
        raise
    except (UnidentifiedImageError, OSError, SyntaxError, ValueError) as exc:
        raise ClientFileError("손상되었거나 지원하지 않는 이미지 파일입니다.") from exc


def client_file_error_response(exc: ClientFileError):
    status = 413 if isinstance(exc, UploadTooLargeError) else 422
    return jsonify({"error": str(exc)}), status


@app.errorhandler(RequestEntityTooLarge)
def handle_request_too_large(_exc):
    if request.path.startswith("/api/"):
        return jsonify({"error": "파일이 너무 큽니다. 100MB 이하 파일만 지원합니다."}), 413
    return "파일이 너무 큽니다. 100MB 이하 파일만 지원합니다.", 413


def human_size(num_bytes: int) -> str:
    size = float(num_bytes)
    for unit in ["B", "KB", "MB", "GB"]:
        if size < 1024:
            return f"{size:.1f}{unit}" if unit != "B" else f"{int(size)}{unit}"
        size /= 1024
    return f"{size:.1f}TB"


# ---------------------------------------------------------------------------
# PDF 압축 (Ghostscript)
# ---------------------------------------------------------------------------
PDF_QUALITY_PRESETS = {
    "low": "/screen",
    "medium": "/ebook",
    "high": "/printer",
}


def compress_pdf(input_path: Path, output_path: Path, quality: str) -> None:
    gs_setting = PDF_QUALITY_PRESETS.get(quality, "/ebook")
    cmd = [
        "gs",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={gs_setting}",
        "-dNOPAUSE",
        "-dQUIET",
        "-dBATCH",
        "-dSAFER",
        f"-sOutputFile={output_path}",
        str(input_path),
    ]
    result = subprocess.run(cmd, capture_output=True, timeout=120)
    if result.returncode != 0 or not output_path.exists():
        raise RuntimeError(
            f"Ghostscript 압축 실패: {result.stderr.decode(errors='ignore')}"
        )


# ---------------------------------------------------------------------------
# 이미지 압축 (Pillow)
# ---------------------------------------------------------------------------
IMG_ALLOWED_EXT = {"jpg", "jpeg", "png", "webp"}

IMG_QUALITY_PRESETS = {
    "high": 90,
    "medium": 75,
    "low": 50,
}

IMG_RESIZE_PRESETS = {
    "100": 1.0,
    "75": 0.75,
    "50": 0.5,
    "25": 0.25,
}

# SNS·용도별 규격 프리셋 (가로, 세로) — 프론트엔드 버튼과 1:1로 매칭됨
IMG_DIMENSION_PRESETS = {
    "instagram_post": (1080, 1080),
    "instagram_story": (1080, 1920),
    "youtube_thumbnail": (1280, 720),
    "blog_featured": (1200, 630),
    "id_photo": (413, 531),
}


def allowed_image_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in IMG_ALLOWED_EXT


def process_image(input_path: Path, output_path: Path, quality: str, resize: str):
    img = Image.open(input_path)
    img = ImageOps.exif_transpose(img)
    orig_w, orig_h = img.size

    scale = IMG_RESIZE_PRESETS.get(resize, 1.0)
    if scale < 1.0:
        new_w = max(1, round(orig_w * scale))
        new_h = max(1, round(orig_h * scale))
        img = img.resize((new_w, new_h), Image.LANCZOS)
    else:
        new_w, new_h = orig_w, orig_h

    q = IMG_QUALITY_PRESETS.get(quality, 75)
    ext = input_path.suffix.lower()

    if ext in (".jpg", ".jpeg"):
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(output_path, "JPEG", quality=q, optimize=True, progressive=True)
    elif ext == ".webp":
        img.save(output_path, "WEBP", quality=q, method=6)
    elif ext == ".png":
        if quality == "low" and img.mode in ("RGB", "RGBA"):
            img = img.convert("RGBA" if "A" in img.mode else "P", palette=Image.ADAPTIVE, colors=256)
        img.save(output_path, "PNG", optimize=True, compress_level=9)
    else:
        img.save(output_path)

    return orig_w, orig_h, new_w, new_h


def resize_to_dimensions(img, target_w, target_h, fit="contain"):
    """target_w/target_h 중 하나 또는 둘 다 지정해 이미지 크기를 맞춘다.
    fit="cover"면 지정한 크기를 정확히 채우도록 가운데를 기준으로 잘라낸다(예: SNS 규격, 증명사진).
    fit="contain"이면 비율을 유지한 채 지정한 크기 안에 들어가도록 축소만 한다(잘림 없음)."""
    orig_w, orig_h = img.size
    target_w = target_w or 0
    target_h = target_h or 0

    if fit == "cover" and target_w and target_h:
        scale = max(target_w / orig_w, target_h / orig_h)
        scaled_w = max(1, round(orig_w * scale))
        scaled_h = max(1, round(orig_h * scale))
        img = img.resize((scaled_w, scaled_h), Image.LANCZOS)
        left = max(0, (scaled_w - target_w) // 2)
        top = max(0, (scaled_h - target_h) // 2)
        img = img.crop((left, top, left + target_w, top + target_h))
        return img, target_w, target_h

    if target_w and target_h:
        scale = min(target_w / orig_w, target_h / orig_h)
    elif target_w:
        scale = target_w / orig_w
    elif target_h:
        scale = target_h / orig_h
    else:
        scale = 1.0

    new_w = max(1, round(orig_w * scale))
    new_h = max(1, round(orig_h * scale))
    img = img.resize((new_w, new_h), Image.LANCZOS)
    return img, new_w, new_h


def _encode_image(img, ext, quality=None):
    buf = io.BytesIO()
    if ext in (".jpg", ".jpeg"):
        save_img = img.convert("RGB") if img.mode in ("RGBA", "P") else img
        save_img.save(buf, "JPEG", quality=quality or 85, optimize=True, progressive=True)
    elif ext == ".webp":
        img.save(buf, "WEBP", quality=quality or 85, method=6)
    elif ext == ".png":
        img.save(buf, "PNG", optimize=True, compress_level=9)
    else:
        img.save(buf, format=(img.format or "JPEG"))
    return buf.getvalue()


def compress_to_target_bytes(img, ext, target_bytes):
    """목표 용량(bytes) 이하가 될 때까지 화질을 낮추고, 그래도 부족하면 크기를 단계적으로
    줄여가며 재인코딩한다. (data, width, height) 를 반환."""
    current = img
    if ext == ".png":
        # PNG은 무손실이라 화질 옵션이 없으므로 크기를 줄여서 용량을 맞춘다.
        data = _encode_image(current, ext)
        for _ in range(12):
            if len(data) <= target_bytes:
                break
            w, h = current.size
            if w <= 80 or h <= 80:
                break
            current = current.resize((max(1, round(w * 0.85)), max(1, round(h * 0.85))), Image.LANCZOS)
            data = _encode_image(current, ext)
        return data, current.size[0], current.size[1]

    # JPEG / WEBP: 화질을 이진 탐색하고, 최저 화질에서도 넘치면 크기를 줄여 재시도한다.
    for _ in range(10):
        lo, hi, best = 10, 95, None
        while lo <= hi:
            mid = (lo + hi) // 2
            data = _encode_image(current, ext, quality=mid)
            if len(data) <= target_bytes:
                best = data
                lo = mid + 1
            else:
                hi = mid - 1
        if best is not None:
            return best, current.size[0], current.size[1]
        w, h = current.size
        if w <= 80 or h <= 80:
            return _encode_image(current, ext, quality=10), w, h
        current = current.resize((max(1, round(w * 0.85)), max(1, round(h * 0.85))), Image.LANCZOS)
    return _encode_image(current, ext, quality=10), current.size[0], current.size[1]


# ---------------------------------------------------------------------------
# PDF → Word 변환 (pdf2docx)
# ---------------------------------------------------------------------------
def convert_pdf_to_docx(input_path: Path, output_path: Path) -> None:
    cv = Converter(str(input_path))
    try:
        cv.convert(str(output_path))
    finally:
        cv.close()


# ---------------------------------------------------------------------------
# PDF 병합 · 분할 (pypdf)
# ---------------------------------------------------------------------------
def merge_pdfs(input_paths, output_path: Path) -> None:
    writer = PdfWriter()
    for p in input_paths:
        reader = PdfReader(str(p))
        if reader.is_encrypted:
            raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
        for page in reader.pages:
            writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)


def parse_page_ranges(range_str: str, total_pages: int) -> list:
    pages = []
    parts = [p.strip() for p in range_str.split(",") if p.strip()]
    if not parts:
        raise ValueError("페이지 범위를 입력해주세요. 예: 1-3, 5, 7-9")
    for part in parts:
        if "-" in part:
            a_str, b_str = part.split("-", 1)
            try:
                a, b = int(a_str), int(b_str)
            except ValueError:
                raise ValueError(f"잘못된 페이지 범위입니다: {part}")
            if a < 1 or b > total_pages or a > b:
                raise ValueError(f"잘못된 페이지 범위입니다: {part} (전체 {total_pages}페이지)")
            pages.extend(range(a - 1, b))
        else:
            try:
                n = int(part)
            except ValueError:
                raise ValueError(f"잘못된 페이지 번호입니다: {part}")
            if n < 1 or n > total_pages:
                raise ValueError(f"잘못된 페이지 번호입니다: {part} (전체 {total_pages}페이지)")
            pages.append(n - 1)
    return pages


def split_pdf_range(input_path: Path, output_path: Path, range_str: str) -> None:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    total = len(reader.pages)
    pages = parse_page_ranges(range_str, total)
    writer = PdfWriter()
    for i in pages:
        writer.add_page(reader.pages[i])
    with open(output_path, "wb") as f:
        writer.write(f)


def split_pdf_individual(input_path: Path, output_zip_path: Path) -> int:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    total = len(reader.pages)
    with zipfile.ZipFile(output_zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for i in range(total):
            writer = PdfWriter()
            writer.add_page(reader.pages[i])
            buf = io.BytesIO()
            writer.write(buf)
            buf.seek(0)
            zf.writestr(f"page_{i + 1:03d}.pdf", buf.read())
    return total


# ---------------------------------------------------------------------------
# 워터마크 추가
# ---------------------------------------------------------------------------
def add_watermark(input_path: Path, output_path: Path, text: str) -> None:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    writer = PdfWriter()

    for page in reader.pages:
        pw = float(page.mediabox.width)
        ph = float(page.mediabox.height)

        buf = io.BytesIO()
        c = pdf_canvas.Canvas(buf, pagesize=(pw, ph))
        c.saveState()
        c.setFillColor(Color(0.5, 0.5, 0.5, alpha=0.35))
        c.setFont("Helvetica-Bold", 40)
        c.translate(pw / 2, ph / 2)
        c.rotate(45)
        step_x, step_y = 320, 220
        for i in range(-3, 4):
            for j in range(-4, 5):
                c.drawCentredString(i * step_x, j * step_y, text)
        c.restoreState()
        c.save()
        buf.seek(0)

        overlay_reader = PdfReader(buf)
        page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# PDF 페이지 회전
# ---------------------------------------------------------------------------
def rotate_pdf(input_path: Path, output_path: Path, angle: int) -> None:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# PDF 비밀번호 설정 · 해제
# ---------------------------------------------------------------------------
def encrypt_pdf(input_path: Path, output_path: Path, password: str) -> None:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("이미 암호가 걸린 PDF입니다.")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    with open(output_path, "wb") as f:
        writer.write(f)


def decrypt_pdf(input_path: Path, output_path: Path, password: str) -> None:
    reader = PdfReader(str(input_path))
    if not reader.is_encrypted:
        raise RuntimeError("암호가 걸려있지 않은 PDF입니다.")
    if reader.decrypt(password) == 0:
        raise RuntimeError("비밀번호가 올바르지 않습니다.")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# QR코드 생성
# ---------------------------------------------------------------------------
def generate_qr(text: str) -> bytes:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=10,
        border=4,
    )
    qr.add_data(text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# 이미지 포맷 변환
# ---------------------------------------------------------------------------
CONVERT_ALLOWED_EXT = {"jpg", "jpeg", "png", "webp", "heic", "heif", "bmp", "tiff", "gif"}
CONVERT_TARGET_FORMATS = {
    "jpg": ("JPEG", ".jpg"),
    "png": ("PNG", ".png"),
    "webp": ("WEBP", ".webp"),
    "bmp": ("BMP", ".bmp"),
    "tiff": ("TIFF", ".tiff"),
}


def allowed_convert_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in CONVERT_ALLOWED_EXT


def convert_image_format(input_path: Path, output_path: Path, target_format: str) -> None:
    fmt_name, _ = CONVERT_TARGET_FORMATS[target_format]
    with PILImage.open(input_path) as source:
        img = ImageOps.exif_transpose(source)
        if fmt_name == "JPEG":
            has_alpha = img.mode in ("RGBA", "LA") or (
                img.mode == "P" and "transparency" in img.info
            )
            if has_alpha:
                rgba = img.convert("RGBA")
                background = PILImage.new("RGB", rgba.size, "white")
                background.paste(rgba, mask=rgba.getchannel("A"))
                img = background
            elif img.mode not in ("RGB", "CMYK"):
                img = img.convert("RGB")
        if fmt_name == "BMP" and img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        img.save(output_path, fmt_name)


# ---------------------------------------------------------------------------
# 전자서명 삽입
# ---------------------------------------------------------------------------
SIGN_POSITIONS = {
    "bottom-right": (0.62, 0.06),
    "bottom-left": (0.06, 0.06),
    "top-right": (0.62, 0.82),
    "top-left": (0.06, 0.82),
    "center": (0.35, 0.45),
}


def apply_signature(input_path: Path, sig_path: Path, output_path: Path, target_page: str, position: str) -> None:
    reader = PdfReader(str(input_path))
    if reader.is_encrypted:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    total = len(reader.pages)

    if target_page == "last":
        target_indices = {total - 1}
    elif target_page == "all":
        target_indices = set(range(total))
    else:
        try:
            n = int(target_page)
        except ValueError:
            raise ValueError("페이지 번호가 올바르지 않습니다.")
        if n < 1 or n > total:
            raise ValueError(f"페이지 번호가 올바르지 않습니다. (전체 {total}페이지)")
        target_indices = {n - 1}

    sig_img = PILImage.open(sig_path)
    sig_w_px, sig_h_px = sig_img.size
    aspect = sig_h_px / sig_w_px

    rel_x, rel_y = SIGN_POSITIONS.get(position, SIGN_POSITIONS["bottom-right"])

    writer = PdfWriter()
    for idx, page in enumerate(reader.pages):
        if idx in target_indices:
            pw = float(page.mediabox.width)
            ph = float(page.mediabox.height)
            sig_w = pw * 0.28
            sig_h = sig_w * aspect

            buf = io.BytesIO()
            c = pdf_canvas.Canvas(buf, pagesize=(pw, ph))
            c.drawImage(
                str(sig_path),
                pw * rel_x,
                ph * rel_y,
                width=sig_w,
                height=sig_h,
                mask="auto",
                preserveAspectRatio=True,
            )
            c.save()
            buf.seek(0)
            overlay_reader = PdfReader(buf)
            page.merge_page(overlay_reader.pages[0])
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)


# ---------------------------------------------------------------------------
# OCR (Tesseract)
# ---------------------------------------------------------------------------
OCR_LANG_MAP = {
    "ko": "kor",
    "en": "eng",
    "ko+en": "kor+eng",
}


def run_tesseract(img, lang: str) -> str:
    try:
        return pytesseract.image_to_string(
            img,
            lang=OCR_LANG_MAP.get(lang, "kor+eng"),
            timeout=OCR_PAGE_TIMEOUT_SECONDS,
        )
    except RuntimeError as exc:
        if "timeout" in str(exc).lower():
            raise ClientFileError(
                "텍스트 인식 시간이 너무 오래 걸립니다. 해상도를 낮추거나 페이지를 나눠 다시 시도해주세요."
            ) from exc
        raise


def ocr_image_file(input_path: Path, lang: str) -> str:
    with PILImage.open(input_path) as img:
        return run_tesseract(img, lang)


def ocr_pdf_file(input_path: Path, lang: str, max_pages: int = 20) -> str:
    texts = []
    with fitz.open(str(input_path)) as doc:
        if doc.needs_pass:
            raise ClientFileError("암호가 걸린 PDF는 지원하지 않습니다.")
        page_count = min(len(doc), max_pages)
        truncated = len(doc) > max_pages
        for i in range(page_count):
            page = doc[i]
            pixel_width = round(page.rect.width / 72 * 200)
            pixel_height = round(page.rect.height / 72 * 200)
            if pixel_width * pixel_height > OCR_PDF_MAX_PAGE_PIXELS:
                raise ClientFileError(
                    "PDF 페이지 크기가 너무 큽니다. 페이지 크기나 해상도를 줄여 다시 시도해주세요."
                )
            pix = page.get_pixmap(dpi=200)
            img = PILImage.frombytes("RGB", (pix.width, pix.height), pix.samples)
            text = run_tesseract(img, lang)
            texts.append(f"--- {i + 1}페이지 ---\n{text.strip()}")
    result = "\n\n".join(texts)
    if truncated:
        result += f"\n\n(처음 {max_pages}페이지까지만 처리했습니다.)"
    return result


# ---------------------------------------------------------------------------
# PDF → PPT 변환
# ---------------------------------------------------------------------------
def convert_pdf_to_pptx(input_path: Path, output_path: Path, max_pages: int = 60) -> None:
    doc = fitz.open(str(input_path))
    if doc.needs_pass:
        raise RuntimeError("암호가 걸린 PDF는 지원하지 않습니다.")
    if len(doc) == 0:
        raise RuntimeError("페이지가 없는 PDF입니다.")

    page0 = doc[0]
    pw_pt, ph_pt = page0.rect.width, page0.rect.height

    prs = Presentation()
    prs.slide_width = Emu(int(pw_pt * 12700))
    prs.slide_height = Emu(int(ph_pt * 12700))
    blank_layout = prs.slide_layouts[6]

    page_count = min(len(doc), max_pages)
    for i in range(page_count):
        page = doc[i]
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    doc.close()
    prs.save(str(output_path))


# ---------------------------------------------------------------------------
# PDF → Excel 변환 (표 추출)
# ---------------------------------------------------------------------------
def convert_pdf_to_xlsx(input_path: Path, output_path: Path, max_pages: int = 30) -> bool:
    wb = Workbook()
    wb.remove(wb.active)
    found_table = False

    with pdfplumber.open(str(input_path)) as pdf:
        if len(pdf.pages) == 0:
            raise RuntimeError("페이지가 없는 PDF입니다.")
        page_count = min(len(pdf.pages), max_pages)
        for i in range(page_count):
            page = pdf.pages[i]
            tables = page.extract_tables()
            if tables:
                for t_idx, table in enumerate(tables):
                    found_table = True
                    sheet_name = f"p{i + 1}_t{t_idx + 1}"[:31]
                    ws = wb.create_sheet(title=sheet_name)
                    for row in table:
                        ws.append(["" if c is None else c for c in row])
            else:
                text = page.extract_text() or ""
                if text.strip():
                    ws = wb.create_sheet(title=f"p{i + 1}_text"[:31])
                    for line in text.splitlines():
                        ws.append([line])

    if len(wb.sheetnames) == 0:
        ws = wb.create_sheet(title="empty")
        ws.append(["추출된 내용이 없습니다."])

    wb.save(str(output_path))
    return found_table


# ---------------------------------------------------------------------------
# 페이지 라우트
# ---------------------------------------------------------------------------
@app.route("/")
def home():
    return render_template(
        "home.html",
        page="home",
        tools=TOOLS,
        categories=TOOL_CATEGORIES,
        site_name=SITE_NAME,
    )


@app.route("/pdf-compress")
def pdf_compress_page():
    return render_template("pdf_compress.html", page="pdf-compress", site_name=SITE_NAME)


@app.route("/image-compress")
def image_compress_page():
    return render_template("image_compress.html", page="image-compress", site_name=SITE_NAME)


@app.route("/pdf-to-word")
def pdf_to_word_page():
    return render_template("pdf_to_word.html", page="pdf-to-word", site_name=SITE_NAME)


@app.route("/pdf-merge-split")
def pdf_merge_split_page():
    return render_template("pdf_merge_split.html", page="pdf-merge-split", site_name=SITE_NAME)


@app.route("/watermark")
def watermark_page():
    return render_template("watermark.html", page="watermark", site_name=SITE_NAME)


@app.route("/pdf-rotate")
def pdf_rotate_page():
    return render_template("pdf_rotate.html", page="pdf-rotate", site_name=SITE_NAME)


@app.route("/pdf-password")
def pdf_password_page():
    return render_template("pdf_password.html", page="pdf-password", site_name=SITE_NAME)


@app.route("/char-counter")
def char_counter_page():
    return render_template("char_counter.html", page="char-counter", site_name=SITE_NAME)


@app.route("/qr-code")
def qr_code_page():
    return render_template("qr_code.html", page="qr-code", site_name=SITE_NAME)


@app.route("/text-diff")
def text_diff_page():
    return render_template("text_diff.html", page="text-diff", site_name=SITE_NAME)


@app.route("/salary-calculator")
def salary_calculator_page():
    return render_template("salary_calculator.html", page="salary-calculator", site_name=SITE_NAME)


@app.route("/insurance-calculator")
def insurance_calculator_page():
    return render_template("insurance_calculator.html", page="insurance-calculator", site_name=SITE_NAME)


@app.route("/severance-calculator")
def severance_calculator_page():
    return render_template("severance_calculator.html", page="severance-calculator", site_name=SITE_NAME)


@app.route("/jeonse-rent-calculator")
def jeonse_rent_calculator_page():
    return render_template("jeonse_rent_calculator.html", page="jeonse-rent-calculator", site_name=SITE_NAME)


@app.route("/image-convert")
def image_convert_page():
    return render_template("image_convert.html", page="image-convert", site_name=SITE_NAME)


@app.route("/pdf-sign")
def pdf_sign_page():
    return render_template("pdf_sign.html", page="pdf-sign", site_name=SITE_NAME)


@app.route("/ocr")
def ocr_page():
    return render_template("ocr.html", page="ocr", site_name=SITE_NAME)


@app.route("/pdf-to-ppt")
def pdf_to_ppt_page():
    return render_template("pdf_to_ppt.html", page="pdf-to-ppt", site_name=SITE_NAME)


@app.route("/pdf-to-excel")
def pdf_to_excel_page():
    return render_template("pdf_to_excel.html", page="pdf-to-excel", site_name=SITE_NAME)


@app.route("/quote-statement")
def quote_statement_page():
    return render_template("quote_statement.html", page="quote-statement", site_name=SITE_NAME)


@app.route("/about")
def about():
    return render_template("about.html", page="about", site_name=SITE_NAME)


@app.route("/contact")
def contact():
    return render_template("contact.html", page="contact", site_name=SITE_NAME)


@app.route("/privacy")
def privacy():
    return render_template("privacy.html", page="privacy", site_name=SITE_NAME)


@app.route("/terms")
def terms():
    return render_template("terms.html", page="terms", site_name=SITE_NAME)


@app.route("/robots.txt")
def robots():
    base = request.host_url.rstrip("/")
    body = "User-agent: *\nAllow: /\n\n" f"Sitemap: {base}/sitemap.xml\n"
    return Response(body, mimetype="text/plain")


@app.route("/sitemap.xml")
def sitemap():
    base = request.host_url.rstrip("/")
    pages = ["", "about", "contact", "privacy", "terms"] + [
        t["slug"] for t in TOOLS if t["available"]
    ]
    entries = []
    for page in pages:
        loc = f"{base}/{page}" if page else f"{base}/"
        if page == "":
            priority = "1.0"
        elif page in ("about", "contact", "privacy", "terms"):
            priority = "0.3"
        else:
            priority = "0.8"
        entries.append(
            f"<url><loc>{loc}</loc><lastmod>{SITE_LAST_UPDATED}</lastmod>"
            f"<changefreq>weekly</changefreq><priority>{priority}</priority></url>"
        )
    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">\n'
        + "\n".join(entries)
        + "\n</urlset>"
    )
    return Response(xml, mimetype="application/xml")


@app.route("/ads.txt")
def ads():
    return app.send_static_file("ads.txt")


@app.route("/manifest.json")
def manifest():
    return app.send_static_file("manifest.json")


@app.route("/service-worker.js")
def service_worker():
    return app.send_static_file("service-worker.js")


@app.route("/favicon.ico")
def favicon():
    return app.send_static_file("favicon.ico")


# ---------------------------------------------------------------------------
# API - PDF 압축
# ---------------------------------------------------------------------------
@app.route("/api/pdf-compress/compress", methods=["POST"])
def api_pdf_compress():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400

    file = request.files["file"]
    quality = request.form.get("quality", "medium")

    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400

    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    if quality not in PDF_QUALITY_PRESETS:
        quality = "medium"

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"

    try:
        original_size = save_upload_limited(
            file, input_path, PDF_MAX_FILE_BYTES, "PDF"
        )
        validate_pdf_file(input_path)
        compress_pdf(input_path, output_path, quality)
    except ClientFileError as exc:
        input_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        input_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "PDF 압축 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500

    compressed_size = output_path.stat().st_size

    used_original = False
    if compressed_size >= original_size:
        output_path.unlink(missing_ok=True)
        output_path = input_path
        compressed_size = original_size
        used_original = True

    download_name = safe_name.rsplit(".", 1)[0] + "_compressed.pdf"

    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/pdf-compress/download/{job_id}?name={download_name}",
            "original_size": original_size,
            "compressed_size": compressed_size,
            "original_size_human": human_size(original_size),
            "compressed_size_human": human_size(compressed_size),
            "ratio": round((1 - compressed_size / original_size) * 100, 1)
            if original_size and not used_original
            else 0,
            "used_original": used_original,
        }
    )


@app.route("/api/pdf-compress/download/<job_id>")
def api_pdf_download(job_id):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.pdf"
    in_path = UPLOAD_DIR / f"{safe_job_id}_in.pdf"

    target = out_path if out_path.exists() else in_path
    if not target.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    download_name = request.args.get("name", "compressed.pdf")

    @after_this_request
    def cleanup(response):
        try:
            in_path.unlink(missing_ok=True)
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        target,
        as_attachment=True,
        download_name=download_name,
        mimetype="application/pdf",
    )


# ---------------------------------------------------------------------------
# API - 이미지 압축
# ---------------------------------------------------------------------------
@app.route("/api/image-compress/process", methods=["POST"])
def api_image_process():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400

    file = request.files["file"]
    mode = request.form.get("mode", "quality")
    if mode not in ("quality", "target_size", "dimensions"):
        mode = "quality"

    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400

    if not allowed_image_file(file.filename):
        return jsonify({"error": "JPG, PNG, WEBP 파일만 업로드할 수 있습니다."}), 400

    # 모드별 파라미터를 먼저 검증한다.
    quality = request.form.get("quality", "medium")
    resize = request.form.get("resize", "100")
    if quality not in IMG_QUALITY_PRESETS:
        quality = "medium"
    if resize not in IMG_RESIZE_PRESETS:
        resize = "100"

    target_bytes = None
    if mode == "target_size":
        target_kb = request.form.get("target_kb", type=float)
        if not target_kb or target_kb <= 0:
            return jsonify({"error": "목표 용량을 입력해주세요."}), 400
        if target_kb > 50 * 1024:
            return jsonify({"error": "목표 용량은 50MB 이하로 입력해주세요."}), 400
        target_bytes = int(target_kb * 1024)

    target_w = target_h = None
    fit = "contain"
    if mode == "dimensions":
        preset = request.form.get("preset", "")
        if preset in IMG_DIMENSION_PRESETS:
            target_w, target_h = IMG_DIMENSION_PRESETS[preset]
            fit = "cover"
        else:
            target_w = request.form.get("target_w", type=int)
            target_h = request.form.get("target_h", type=int)
            fit = "cover" if request.form.get("fit") == "cover" else "contain"
        if not target_w and not target_h:
            return jsonify({"error": "가로 또는 세로 크기를 입력해주세요."}), 400
        if (target_w and (target_w <= 0 or target_w > 8000)) or (
            target_h and (target_h <= 0 or target_h > 8000)
        ):
            return jsonify({"error": "크기는 1~8000px 사이로 입력해주세요."}), 400
        if fit == "cover" and not (target_w and target_h):
            return jsonify({"error": "정확한 크기로 맞추려면 가로·세로를 모두 입력해주세요."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "image.jpg"
    ext = Path(safe_name).suffix.lower() or ".jpg"
    input_path = UPLOAD_DIR / f"{job_id}_in{ext}"
    output_path = UPLOAD_DIR / f"{job_id}_out{ext}"

    try:
        original_size = save_upload_limited(
            file, input_path, IMAGE_COMPRESS_MAX_FILE_BYTES, "이미지"
        )
        validate_image_file(input_path, max_pixels=IMAGE_COMPRESS_MAX_PIXELS)
        if mode == "target_size":
            img = Image.open(input_path)
            img = ImageOps.exif_transpose(img)
            orig_w, orig_h = img.size
            data, new_w, new_h = compress_to_target_bytes(img, ext, target_bytes)
            output_path.write_bytes(data)
        elif mode == "dimensions":
            img = Image.open(input_path)
            img = ImageOps.exif_transpose(img)
            orig_w, orig_h = img.size
            resized, new_w, new_h = resize_to_dimensions(img, target_w, target_h, fit)
            q = IMG_QUALITY_PRESETS.get(quality, 85)
            data = _encode_image(resized, ext, quality=q)
            output_path.write_bytes(data)
        else:
            orig_w, orig_h, new_w, new_h = process_image(input_path, output_path, quality, resize)
    except ClientFileError as exc:
        input_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        input_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "이미지 처리 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500

    compressed_size = output_path.stat().st_size

    used_original = False
    if mode == "quality" and compressed_size >= original_size and resize == "100":
        output_path.unlink(missing_ok=True)
        output_path = input_path
        compressed_size = original_size
        used_original = True

    stem = safe_name.rsplit(".", 1)[0]
    download_name = f"{stem}_compressed{ext}"

    response = {
        "job_id": job_id,
        "download_url": f"/api/image-compress/download/{job_id}?name={download_name}&ext={ext}",
        "original_size": original_size,
        "compressed_size": compressed_size,
        "original_size_human": human_size(original_size),
        "compressed_size_human": human_size(compressed_size),
        "ratio": round((1 - compressed_size / original_size) * 100, 1)
        if original_size and not used_original
        else 0,
        "original_dimensions": f"{orig_w}×{orig_h}",
        "new_dimensions": f"{new_w}×{new_h}",
        "used_original": used_original,
        "mode": mode,
    }
    if mode == "target_size":
        response["target_reached"] = compressed_size <= target_bytes
    return jsonify(response)


@app.route("/api/image-compress/download/<job_id>")
def api_image_download(job_id):
    safe_job_id = secure_filename(job_id)
    ext = request.args.get("ext", ".jpg")
    if ext not in (".jpg", ".jpeg", ".png", ".webp"):
        ext = ".jpg"

    out_path = UPLOAD_DIR / f"{safe_job_id}_out{ext}"
    in_path = UPLOAD_DIR / f"{safe_job_id}_in{ext}"

    target = out_path if out_path.exists() else in_path
    if not target.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    download_name = request.args.get("name", f"compressed{ext}")
    mimetype = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".webp": "image/webp",
    }.get(ext, "application/octet-stream")

    @after_this_request
    def cleanup(response):
        try:
            in_path.unlink(missing_ok=True)
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        target,
        as_attachment=True,
        download_name=download_name,
        mimetype=mimetype,
    )


# ---------------------------------------------------------------------------
# API - PDF → Word 변환
# ---------------------------------------------------------------------------
@app.route("/api/pdf-to-word/convert", methods=["POST"])
def api_pdf_to_word():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.docx"

    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        convert_pdf_to_docx(input_path, output_path)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "Word 변환에 실패했습니다. 텍스트 기반 PDF인지 확인해주세요."}), 422
    finally:
        input_path.unlink(missing_ok=True)

    if not output_path.exists():
        return jsonify({"error": "변환에 실패했습니다."}), 500

    download_name = safe_name.rsplit(".", 1)[0] + ".docx"

    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/pdf-to-word/download/{job_id}?name={download_name}",
        }
    )


@app.route("/api/pdf-to-word/download/<job_id>")
def api_pdf_to_word_download(job_id):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.docx"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    download_name = request.args.get("name", "converted.docx")

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name=download_name,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


# ---------------------------------------------------------------------------
# API - PDF 병합 · 분할
# ---------------------------------------------------------------------------
@app.route("/api/pdf-merge-split/merge", methods=["POST"])
def api_pdf_merge():
    files = [f for f in request.files.getlist("files") if f and f.filename]
    if len(files) < 2:
        return jsonify({"error": "합칠 PDF 파일을 2개 이상 선택해주세요."}), 400
    for f in files:
        if not f.filename.lower().endswith(".pdf"):
            return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    input_paths = []
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"

    try:
        remaining_bytes = PDF_MAX_FILE_BYTES
        for i, f in enumerate(files):
            p = UPLOAD_DIR / f"{job_id}_in{i}.pdf"
            saved_bytes = save_upload_limited(f, p, remaining_bytes, "전체 PDF")
            remaining_bytes -= saved_bytes
            input_paths.append(p)
            validate_pdf_file(p)
        merge_pdfs(input_paths, output_path)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "PDF 병합 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500
    finally:
        for p in input_paths:
            p.unlink(missing_ok=True)

    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/pdf-merge-split/download/{job_id}?name=merged.pdf&type=pdf",
        }
    )


@app.route("/api/pdf-merge-split/split", methods=["POST"])
def api_pdf_split():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400

    file = request.files["file"]
    mode = request.form.get("mode", "individual")
    page_range = request.form.get("range", "")

    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"

    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        if mode == "range":
            output_path = UPLOAD_DIR / f"{job_id}_out.pdf"
            split_pdf_range(input_path, output_path, page_range)
            download_name = safe_name.rsplit(".", 1)[0] + "_split.pdf"
            file_type = "pdf"
        else:
            output_path = UPLOAD_DIR / f"{job_id}_out.zip"
            split_pdf_individual(input_path, output_path)
            download_name = safe_name.rsplit(".", 1)[0] + "_pages.zip"
            file_type = "zip"
    except ClientFileError as exc:
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": "PDF 분할 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/pdf-merge-split/download/{job_id}?name={download_name}&type={file_type}",
        }
    )


@app.route("/api/pdf-merge-split/download/<job_id>")
def api_pdf_merge_split_download(job_id):
    safe_job_id = secure_filename(job_id)
    file_type = request.args.get("type", "pdf")
    ext = ".zip" if file_type == "zip" else ".pdf"
    out_path = UPLOAD_DIR / f"{safe_job_id}_out{ext}"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    download_name = request.args.get("name", f"result{ext}")
    mimetype = "application/zip" if file_type == "zip" else "application/pdf"

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name=download_name,
        mimetype=mimetype,
    )


# ---------------------------------------------------------------------------
# API - 워터마크 추가
# ---------------------------------------------------------------------------
@app.route("/api/watermark/add", methods=["POST"])
def api_watermark_add():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    text = (request.form.get("text") or "").strip()
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400
    if not text:
        return jsonify({"error": "워터마크 텍스트를 입력해주세요."}), 400
    if len(text) > 40:
        text = text[:40]

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"
    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        add_watermark(input_path, output_path, text)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "워터마크 처리 중 오류가 발생했습니다."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + "_watermark.pdf"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/watermark/download/{job_id}?name={download_name}"}
    )


@app.route("/api/watermark/download/<job_id>")
def api_watermark_download(job_id):
    return _simple_pdf_download(job_id, "watermarked.pdf")


# ---------------------------------------------------------------------------
# API - PDF 페이지 회전
# ---------------------------------------------------------------------------
@app.route("/api/pdf-rotate/rotate", methods=["POST"])
def api_pdf_rotate():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    try:
        angle = int(request.form.get("angle", "90"))
    except ValueError:
        angle = 90
    if angle not in (90, 180, 270):
        angle = 90
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"
    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        rotate_pdf(input_path, output_path, angle)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "PDF 회전 중 오류가 발생했습니다."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + "_rotated.pdf"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/pdf-rotate/download/{job_id}?name={download_name}"}
    )


@app.route("/api/pdf-rotate/download/<job_id>")
def api_pdf_rotate_download(job_id):
    return _simple_pdf_download(job_id, "rotated.pdf")


# ---------------------------------------------------------------------------
# API - PDF 비밀번호 설정 · 해제
# ---------------------------------------------------------------------------
@app.route("/api/pdf-password/encrypt", methods=["POST"])
def api_pdf_encrypt():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    password = request.form.get("password", "")
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400
    if not password:
        return jsonify({"error": "비밀번호를 입력해주세요."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"
    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        encrypt_pdf(input_path, output_path, password)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "PDF 암호 설정 중 오류가 발생했습니다."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + "_locked.pdf"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/pdf-password/download/{job_id}?name={download_name}"}
    )


@app.route("/api/pdf-password/decrypt", methods=["POST"])
def api_pdf_decrypt():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    password = request.form.get("password", "")
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400
    if not password:
        return jsonify({"error": "비밀번호를 입력해주세요."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"
    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        _page_count, encrypted = validate_pdf_file(input_path, allow_encrypted=True)
        if not encrypted:
            raise ClientFileError("암호가 설정된 PDF 파일을 선택해주세요.")
        decrypt_pdf(input_path, output_path, password)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "비밀번호가 맞지 않거나 PDF를 해제할 수 없습니다."}), 422
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + "_unlocked.pdf"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/pdf-password/download/{job_id}?name={download_name}"}
    )


@app.route("/api/pdf-password/download/<job_id>")
def api_pdf_password_download(job_id):
    return _simple_pdf_download(job_id, "result.pdf")


# ---------------------------------------------------------------------------
# API - QR코드 생성
# ---------------------------------------------------------------------------
@app.route("/api/qr-code/generate", methods=["POST"])
def api_qr_generate():
    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()
    if not text:
        return jsonify({"error": "텍스트나 URL을 입력해주세요."}), 400
    if len(text) > 1000:
        return jsonify({"error": "텍스트가 너무 깁니다. 1000자 이하로 입력해주세요."}), 400

    try:
        png_bytes = generate_qr(text)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": f"QR코드 생성 실패: {exc}"}), 500

    job_id = uuid.uuid4().hex
    out_path = UPLOAD_DIR / f"{job_id}_out.png"
    out_path.write_bytes(png_bytes)

    return jsonify({"job_id": job_id, "download_url": f"/api/qr-code/download/{job_id}"})


@app.route("/api/qr-code/download/<job_id>")
def api_qr_download(job_id):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.png"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(out_path, mimetype="image/png", as_attachment=False, download_name="qrcode.png")


# ---------------------------------------------------------------------------
# API - 이미지 포맷 변환
# ---------------------------------------------------------------------------
@app.route("/api/image-convert/convert", methods=["POST"])
def api_image_convert():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    target = request.form.get("target", "jpg").lower()

    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not allowed_convert_file(file.filename):
        return jsonify({"error": "지원하지 않는 이미지 형식입니다."}), 400
    if target not in CONVERT_TARGET_FORMATS:
        return jsonify({"error": "지원하지 않는 변환 형식입니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "image"
    in_ext = Path(safe_name).suffix.lower() or ".jpg"
    _, out_ext = CONVERT_TARGET_FORMATS[target]
    input_path = UPLOAD_DIR / f"{job_id}_in{in_ext}"
    output_path = UPLOAD_DIR / f"{job_id}_out{out_ext}"

    try:
        save_upload_limited(
            file, input_path, IMAGE_CONVERT_MAX_FILE_BYTES, "이미지"
        )
        validate_image_file(input_path, max_pixels=IMAGE_CONVERT_MAX_PIXELS)
        convert_image_format(input_path, output_path, target)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "이미지 변환 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    stem = safe_name.rsplit(".", 1)[0]
    download_name = f"{stem}{out_ext}"
    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/image-convert/download/{job_id}?name={download_name}&ext={out_ext}",
        }
    )


@app.route("/api/image-convert/download/<job_id>")
def api_image_convert_download(job_id):
    safe_job_id = secure_filename(job_id)
    ext = request.args.get("ext", ".jpg")
    if ext not in (".jpg", ".png", ".webp", ".bmp", ".tiff"):
        ext = ".jpg"
    out_path = UPLOAD_DIR / f"{safe_job_id}_out{ext}"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404

    download_name = request.args.get("name", f"converted{ext}")
    mimetype = {
        ".jpg": "image/jpeg", ".png": "image/png", ".webp": "image/webp",
        ".bmp": "image/bmp", ".tiff": "image/tiff",
    }.get(ext, "application/octet-stream")

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(out_path, as_attachment=True, download_name=download_name, mimetype=mimetype)


# ---------------------------------------------------------------------------
# API - 전자서명 삽입
# ---------------------------------------------------------------------------
@app.route("/api/pdf-sign/apply", methods=["POST"])
def api_pdf_sign():
    if "file" not in request.files or "signature" not in request.files:
        return jsonify({"error": "PDF 파일과 서명 이미지를 모두 첨부해주세요."}), 400
    file = request.files["file"]
    sig_file = request.files["signature"]
    target_page = request.form.get("target_page", "last")
    position = request.form.get("position", "bottom-right")

    if file.filename == "" or sig_file.filename == "":
        return jsonify({"error": "PDF 파일과 서명 이미지를 모두 첨부해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    sig_path = UPLOAD_DIR / f"{job_id}_sig.png"
    output_path = UPLOAD_DIR / f"{job_id}_out.pdf"

    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        save_upload_limited(
            sig_file, sig_path, SIGNATURE_MAX_FILE_BYTES, "서명 이미지"
        )
        validate_image_file(sig_path, max_pixels=SIGNATURE_MAX_PIXELS)
        apply_signature(input_path, sig_path, output_path, target_page, position)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "서명 삽입 중 오류가 발생했습니다."}), 500
    finally:
        input_path.unlink(missing_ok=True)
        sig_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + "_signed.pdf"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/pdf-sign/download/{job_id}?name={download_name}"}
    )


@app.route("/api/pdf-sign/download/<job_id>")
def api_pdf_sign_download(job_id):
    return _simple_pdf_download(job_id, "signed.pdf")


# ---------------------------------------------------------------------------
# API - OCR
# ---------------------------------------------------------------------------
@app.route("/api/ocr/extract", methods=["POST"])
def api_ocr_extract():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    lang = request.form.get("lang", "ko+en")
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400

    filename_lower = file.filename.lower()
    is_pdf = filename_lower.endswith(".pdf")
    is_image = any(filename_lower.endswith(ext) for ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"))
    if not (is_pdf or is_image):
        return jsonify({"error": "PDF 또는 이미지 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    ext = ".pdf" if is_pdf else Path(secure_filename(file.filename)).suffix.lower()
    input_path = UPLOAD_DIR / f"{job_id}_in{ext}"

    try:
        save_upload_limited(file, input_path, OCR_MAX_FILE_BYTES, "OCR")
        if is_pdf:
            validate_pdf_file(input_path)
            text = ocr_pdf_file(input_path, lang)
        else:
            validate_image_file(input_path, max_pixels=OCR_IMAGE_MAX_PIXELS)
            text = ocr_image_file(input_path, lang)
    except ClientFileError as exc:
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": "텍스트 인식 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    return jsonify({"text": text})


# ---------------------------------------------------------------------------
# API - PDF → PPT 변환
# ---------------------------------------------------------------------------
@app.route("/api/pdf-to-ppt/convert", methods=["POST"])
def api_pdf_to_ppt():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.pptx"

    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        convert_pdf_to_pptx(input_path, output_path)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "PPT 변환 중 오류가 발생했습니다."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + ".pptx"
    return jsonify(
        {"job_id": job_id, "download_url": f"/api/pdf-to-ppt/download/{job_id}?name={download_name}"}
    )


@app.route("/api/pdf-to-ppt/download/<job_id>")
def api_pdf_to_ppt_download(job_id):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.pptx"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404
    download_name = request.args.get("name", "converted.pptx")

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name=download_name,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ---------------------------------------------------------------------------
# API - PDF → Excel 변환
# ---------------------------------------------------------------------------
@app.route("/api/pdf-to-excel/convert", methods=["POST"])
def api_pdf_to_excel():
    if "file" not in request.files:
        return jsonify({"error": "파일을 선택해주세요."}), 400
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "파일을 선택해주세요."}), 400
    if not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": "PDF 파일만 업로드할 수 있습니다."}), 400

    job_id = uuid.uuid4().hex
    safe_name = secure_filename(file.filename) or "document.pdf"
    input_path = UPLOAD_DIR / f"{job_id}_in.pdf"
    output_path = UPLOAD_DIR / f"{job_id}_out.xlsx"

    try:
        save_upload_limited(file, input_path, PDF_MAX_FILE_BYTES, "PDF")
        validate_pdf_file(input_path)
        found_table = convert_pdf_to_xlsx(input_path, output_path)
    except ClientFileError as exc:
        output_path.unlink(missing_ok=True)
        return client_file_error_response(exc)
    except Exception as exc:  # noqa: BLE001
        output_path.unlink(missing_ok=True)
        return jsonify({"error": "Excel 변환 중 오류가 발생했습니다. 잠시 후 다시 시도해주세요."}), 500
    finally:
        input_path.unlink(missing_ok=True)

    download_name = safe_name.rsplit(".", 1)[0] + ".xlsx"
    return jsonify(
        {
            "job_id": job_id,
            "download_url": f"/api/pdf-to-excel/download/{job_id}?name={download_name}",
            "found_table": found_table,
        }
    )


@app.route("/api/pdf-to-excel/download/<job_id>")
def api_pdf_to_excel_download(job_id):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.xlsx"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404
    download_name = request.args.get("name", "converted.xlsx")

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(
        out_path,
        as_attachment=True,
        download_name=download_name,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def _simple_pdf_download(job_id, default_name):
    safe_job_id = secure_filename(job_id)
    out_path = UPLOAD_DIR / f"{safe_job_id}_out.pdf"
    if not out_path.exists():
        return jsonify({"error": "파일을 찾을 수 없습니다. 다시 시도해주세요."}), 404
    download_name = request.args.get("name", default_name)

    @after_this_request
    def cleanup(response):
        try:
            out_path.unlink(missing_ok=True)
        except Exception:  # noqa: BLE001
            pass
        return response

    return send_file(out_path, as_attachment=True, download_name=download_name, mimetype="application/pdf")


@app.errorhandler(413)
def too_large(_e):
    return jsonify({"error": "파일이 너무 큽니다. 100MB 이하 파일만 지원합니다."}), 413


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5002))
    app.run(host="0.0.0.0", port=port, debug=False)
